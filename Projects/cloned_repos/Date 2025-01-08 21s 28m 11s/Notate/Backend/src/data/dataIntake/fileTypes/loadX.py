import pandas as pd
import json
import markdown
from bs4 import BeautifulSoup
from pptx import Presentation
from langchain_community.document_loaders import Docx2txtLoader
from langchain_community.document_loaders.csv_loader import CSVLoader
import logging


def load_py(file):
    try:
        with open(file, 'r', encoding='utf-8') as f:
            content = f.read()
            return content.strip()
    except Exception as e:
        print(f"Error loading PY: {str(e)}")
        return None


def load_pdf(file):
    try:
        """   loader = MathpixPDFLoader(file)
          docs = loader.load()
          print(docs)
          print(f"Successfully loaded PDF file: {file}")
          print(f"Number of pages: {len(docs)}")
          print(f"PDF metadata: {docs[0].metadata}")
          print(docs[0].page_content)
          return docs[0].page_content """
        return None
    except Exception as e:
        print(f"Error loading PDF: {str(e)}")
        return None


def load_docx(file):
    try:
        loader = Docx2txtLoader(file)
        data = loader.load()
        print(data)
        return data[0].page_content
    except Exception as e:
        print(f"Error loading DOCX: {str(e)}")
        return None


def load_txt(file):
    try:
        with open(file, 'r', encoding='utf-8') as f:
            return f.read().strip()
    except Exception as e:
        print(f"Error loading TXT: {str(e)}")
        return None


def load_md(file):
    try:
        with open(file, 'r', encoding='utf-8') as f:
            md_text = f.read()
            html = markdown.markdown(md_text)
            soup = BeautifulSoup(html, 'html.parser')
            return soup.get_text().strip()
    except Exception as e:
        print(f"Error loading MD: {str(e)}")
        return None


def load_html(file_path: str) -> str:
    """Load and process HTML file content"""
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()

        # Parse HTML with BeautifulSoup
        soup = BeautifulSoup(content, 'html.parser')

        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()

        # Get text content
        text = soup.get_text()

        # Break into lines and remove leading/trailing space
        lines = (line.strip() for line in text.splitlines())

        # Break multi-headlines into a line each
        chunks = (phrase.strip()
                  for line in lines for phrase in line.split("  "))

        # Drop blank lines
        text = ' '.join(chunk for chunk in chunks if chunk)

        return text
    except Exception as e:
        logging.error(f"Error loading HTML file {file_path}: {str(e)}")
        return None


def load_csv(file):
    try:
        loader = CSVLoader(file)
        data = loader.load()
        return data
    except Exception as e:
        print(f"Error loading CSV: {str(e)}")
        return None


def load_json(file):
    try:
        with open(file, 'r', encoding='utf-8') as f:
            data = json.load(f)
            return json.dumps(data, indent=2)
    except Exception as e:
        print(f"Error loading JSON: {str(e)}")
        return None


def load_pptx(file):
    try:
        prs = Presentation(file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text).strip()
    except Exception as e:
        print(f"Error loading PPTX: {str(e)}")
        return None


def load_xlsx(file):
    try:
        df = pd.read_excel(file)
        return df.to_string().strip()
    except Exception as e:
        print(f"Error loading XLSX: {str(e)}")
        return None
